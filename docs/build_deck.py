#!/usr/bin/env python3
"""Build docs/deck.pptx from docs/deck.html.

The HTML deck is the source of truth for both the slides and the speaker notes:
each `<section class="slide">` carries an `<aside class="notes">` that this script
lifts into the PowerPoint notes pane. Press `n` in the HTML deck to see the same text.

Each slide is rendered by headless Chrome at 2560x1440 and placed full-bleed, so the
PPTX looks exactly like the deck rather than approximating it in PowerPoint shapes.
That makes the slides images: edit `deck.html` and re-run this, do not edit the pptx.

The HTML follows the viewer's theme; a pptx cannot, so the theme is baked and has to be
chosen. Light is the default because conference projectors wash out dark backgrounds and
this deck leans on hairline rules and muted grey text, both of which go first.

    python docs/build_deck.py                 # light, for a projector
    python docs/build_deck.py --theme dark    # dark, for a screen
"""
from __future__ import annotations

import argparse
import html
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches

DOCS = Path(__file__).resolve().parent
HTML = DOCS / "deck.html"
PPTX = DOCS / "deck.pptx"

WIDTH, HEIGHT = 2560, 1440          # 16:9, comfortably above projector resolution
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

CHROME = (
    "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
    "/Applications/Chromium.app/Contents/MacOS/Chromium",
    "google-chrome",
    "chromium",
)

# Strip the presenter chrome and let one slide fill the viewport exactly.
FULL_BLEED = """
<style>
  #rail, #notes { display: none !important; }
  #deck { padding: 0 !important; height: 100vh !important; }
  .slide { width: 100vw !important; height: 100vh !important; max-width: none !important;
           aspect-ratio: auto !important; border: 0 !important; border-radius: 0 !important;
           box-shadow: none !important; font-size: 1.62vw !important; }
</style>
"""


def page_source(markup: str, theme: str) -> str:
    """Force the theme: the deck's tokens key off a data-theme attribute on the root."""
    return (f'<!doctype html><html data-theme="{theme}">'
            f'<base href="file://{DOCS}/">' + FULL_BLEED + markup + "</html>")


def find_chrome() -> str:
    for c in CHROME:
        if Path(c).exists() or shutil.which(c):
            return c
    sys.exit("no Chrome or Chromium found — install one, or edit CHROME in this script")


def slides_and_notes(markup: str) -> list[str]:
    """One notes string per slide, in order. Fails loudly if a slide has none."""
    sections = re.findall(r'<section class="slide[^"]*">(.*?)</section>', markup, re.S)
    notes: list[str] = []
    for i, sec in enumerate(sections, 1):
        m = re.search(r'<aside class="notes">(.*?)</aside>', sec, re.S)
        if not m:
            sys.exit(f"slide {i} has no <aside class=\"notes\"> block")
        text = re.sub(r"<[^>]+>", "", m.group(1))
        notes.append(html.unescape(text).strip())
    return notes


def shoot(chrome: str, page: Path, n: int, out: Path) -> None:
    subprocess.run(
        [chrome, "--headless", "--disable-gpu", "--hide-scrollbars",
         f"--window-size={WIDTH},{HEIGHT}", f"--screenshot={out}",
         f"file://{page}#{n}"],
        check=True, capture_output=True,
    )
    if not out.exists():
        sys.exit(f"chrome produced no image for slide {n}")


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    ap.add_argument("--theme", choices=("light", "dark"), default="light",
                    help="baked into every slide; light is safer on a projector")
    ap.add_argument("--out", type=Path, default=PPTX)
    args = ap.parse_args()

    if not HTML.exists():
        sys.exit(f"missing {HTML}")
    chrome = find_chrome()
    markup = HTML.read_text(encoding="utf-8")
    notes = slides_and_notes(markup)
    print(f"{len(notes)} slides, {sum(len(n.split()) for n in notes)} words of notes, "
          f"{args.theme} theme")

    deck = Presentation()
    deck.slide_width, deck.slide_height = SLIDE_W, SLIDE_H
    blank = deck.slide_layouts[6]

    with tempfile.TemporaryDirectory() as tmp:
        tmpdir = Path(tmp)
        page = tmpdir / "full_bleed.html"
        page.write_text(page_source(markup, args.theme), encoding="utf-8")

        for n, note in enumerate(notes, 1):
            png = tmpdir / f"slide{n:02d}.png"
            shoot(chrome, page, n, png)
            slide = deck.slides.add_slide(blank)
            slide.shapes.add_picture(str(png), Emu(0), Emu(0),
                                     width=SLIDE_W, height=SLIDE_H)
            slide.notes_slide.notes_text_frame.text = note
            print(f"  slide {n}  {png.stat().st_size // 1024:>5} KB  "
                  f"{len(note.split()):>3} words of notes")

    deck.save(args.out)
    print(f"wrote {args.out}  {args.out.stat().st_size / 1e6:.1f} MB")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
